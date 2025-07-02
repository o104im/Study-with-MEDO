from pptx import Presentation
import nltk
import spacy
from nltk.tokenize import sent_tokenize
import re

# تحميل نموذج اللغة الإنجليزية
nlp = spacy.load('en_core_web_sm')

class PPTXProcessor:
    def __init__(self, file_path):
        self.file_path = file_path
        self.text = self._extract_text()
        
    def _extract_text(self):
        """استخراج النص من ملف PowerPoint"""
        text = ""
        try:
            prs = Presentation(self.file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return ""
    
    def get_text(self):
        """الحصول على النص الكامل"""
        return self.text
    
    def summarize(self, ratio=0.3):
        """تلخيص النص باستخدام خوارزمية بسيطة"""
        if not self.text:
            return ""
        
        # تقسيم النص إلى جمل
        sentences = sent_tokenize(self.text)
        
        # إذا كان عدد الجمل قليلاً، إرجاع النص كاملاً
        if len(sentences) <= 5:
            return self.text
        
        # تحليل النص وحساب تكرار الكلمات
        word_frequencies = {}
        
        # تحليل النص باستخدام spaCy
        doc = nlp(self.text)
        
        # استخراج الكلمات المهمة (بدون كلمات التوقف)
        for token in doc:
            if not token.is_stop and not token.is_punct and not token.is_space:
                word = token.text.lower()
                if word in word_frequencies:
                    word_frequencies[word] += 1
                else:
                    word_frequencies[word] = 1
        
        # تطبيع التكرارات
        max_frequency = max(word_frequencies.values()) if word_frequencies else 1
        for word in word_frequencies:
            word_frequencies[word] = word_frequencies[word] / max_frequency
        
        # حساب درجة أهمية كل جملة
        sentence_scores = {}
        for i, sentence in enumerate(sentences):
            for word in nlp(sentence.lower()):
                if word.text.lower() in word_frequencies:
                    if i in sentence_scores:
                        sentence_scores[i] += word_frequencies[word.text.lower()]
                    else:
                        sentence_scores[i] = word_frequencies[word.text.lower()]
        
        # اختيار الجمل الأكثر أهمية
        num_sentences = max(1, int(len(sentences) * ratio))
        selected_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)[:num_sentences]
        selected_sentences = sorted(selected_sentences, key=lambda x: x[0])
        
        # إنشاء الملخص
        summary = " ".join([sentences[i] for i, _ in selected_sentences])
        
        return summary
    
    def generate_questions(self, question_type='multiple-choice', num_questions=5):
        """توليد أسئلة من النص"""
        if not self.text or len(self.text) < 100:
            return []
        
        # تحليل النص
        doc = nlp(self.text)
        
        # استخراج الجمل المهمة
        important_sentences = []
        for sent in doc.sents:
            # تجاهل الجمل القصيرة جدًا
            if len(sent.text.split()) > 5:
                important_sentences.append(sent.text)
        
        # إذا لم يتم العثور على جمل كافية
        if len(important_sentences) < num_questions:
            num_questions = len(important_sentences)
        
        # اختيار جمل عشوائية لتوليد الأسئلة
        import random
        selected_sentences = random.sample(important_sentences, min(num_questions, len(important_sentences)))
        
        questions = []
        
        if question_type == 'multiple-choice':
            for i, sentence in enumerate(selected_sentences):
                # تحليل الجملة
                sent_doc = nlp(sentence)
                
                # البحث عن كلمات مهمة في الجملة
                important_words = []
                for token in sent_doc:
                    if (token.pos_ in ['NOUN', 'PROPN', 'NUM'] and 
                        not token.is_stop and len(token.text) > 2):
                        important_words.append(token.text)
                
                if not important_words:
                    continue
                
                # اختيار كلمة عشوائية لتكون الإجابة الصحيحة
                answer = random.choice(important_words)
                
                # إنشاء السؤال
                question_text = sentence.replace(answer, "______")
                
                # إنشاء خيارات
                options = [answer]
                
                # إضافة خيارات عشوائية
                all_nouns = [token.text for token in doc if token.pos_ in ['NOUN', 'PROPN'] 
                            and not token.is_stop and len(token.text) > 2 
                            and token.text != answer]
                
                if len(all_nouns) >= 3:
                    options.extend(random.sample(all_nouns, 3))
                else:
                    # إذا لم يكن هناك ما يكفي من الأسماء، أضف بعض الكلمات العشوائية
                    options.extend(all_nouns)
                    while len(options) < 4:
                        options.append(f"Option {len(options)}")
                
                # خلط الخيارات
                random.shuffle(options)
                
                questions.append({
                    'id': i + 1,
                    'question': question_text,
                    'options': options,
                    'correctAnswer': answer
                })
        
        elif question_type == 'short-answer':
            for i, sentence in enumerate(selected_sentences):
                # تحليل الجملة
                sent_doc = nlp(sentence)
                
                # استخراج الكلمات المفتاحية
                keywords = []
                for ent in sent_doc.ents:
                    if ent.label_ in ['PERSON', 'ORG', 'GPE', 'DATE', 'NORP']:
                        keywords.append(ent.text)
                
                # إذا لم يتم العثور على كيانات، ابحث عن أسماء
                if not keywords:
                    for token in sent_doc:
                        if token.pos_ in ['NOUN', 'PROPN'] and not token.is_stop:
                            keywords.append(token.text)
                
                if not keywords:
                    continue
                
                # إنشاء سؤال
                keyword = random.choice(keywords)
                question_text = sentence.replace(keyword, "______")
                question_text = f"ما هو/هي {keyword}؟" if random.random() > 0.5 else f"اشرح مفهوم {keyword}."
                
                questions.append({
                    'id': i + 1,
                    'question': question_text,
                    'answer': sentence
                })
        
        elif question_type == 'true-false':
            for i, sentence in enumerate(selected_sentences):
                # نسخة صحيحة من الجملة
                correct_version = sentence
                
                # إنشاء نسخة خاطئة من الجملة في 50% من الحالات
                is_true = random.random() > 0.5
                
                if not is_true:
                    # تحليل الجملة
                    sent_doc = nlp(sentence)
                    
                    # محاولة تغيير جزء من الجملة
                    for token in sent_doc:
                        if token.pos_ in ['NOUN', 'NUM', 'ADJ']:
                            # تغيير الكلمة بكلمة عشوائية من نفس النوع
                            replacements = [t.text for t in doc if t.pos_ == token.pos_ and t.text != token.text]
                            if replacements:
                                replacement = random.choice(replacements)
                                correct_version = sentence  # الجملة الأصلية هي الصحيحة
                                sentence = sentence.replace(token.text, replacement, 1)
                                break
                
                questions.append({
                    'id': i + 1,
                    'question': sentence,
                    'answer': is_true
                })
        
        return questions
    
    def generate_mindmap_data(self):
        """توليد بيانات لخريطة ذهنية"""
        if not self.text:
            return {}
        
        # تحليل النص
        doc = nlp(self.text)
        
        # استخراج الموضوع الرئيسي (عنوان المستند أو الجملة الأولى)
        main_topic = doc[:10].text if len(doc) > 10 else doc.text
        
        # استخراج الكيانات المهمة
        entities = {}
        for ent in doc.ents:
            if ent.label_ not in entities:
                entities[ent.label_] = []
            if ent.text not in entities[ent.label_]:
                entities[ent.label_].append(ent.text)
        
        # استخراج الأفكار الرئيسية (الجمل المهمة)
        main_ideas = []
        for sent in doc.sents:
            # تجاهل الجمل القصيرة جدًا
            if len(sent.text.split()) > 8:
                main_ideas.append(sent.text)
        
        # اختيار عدد محدود من الأفكار الرئيسية
        import random
        selected_ideas = random.sample(main_ideas, min(4, len(main_ideas))) if len(main_ideas) > 4 else main_ideas
        
        # إنشاء بيانات الخريطة الذهنية
        mindmap_data = {
            'main_topic': main_topic,
            'branches': []
        }
        
        # إضافة فروع من الكيانات المهمة
        for label, items in entities.items():
            if len(items) > 0:
                branch = {
                    'name': label,
                    'children': [{'name': item} for item in items[:3]]  # الحد من عدد العناصر
                }
                mindmap_data['branches'].append(branch)
        
        # إضافة فروع من الأفكار الرئيسية
        for i, idea in enumerate(selected_ideas):
            # تقصير الفكرة إذا كانت طويلة جدًا
            short_idea = idea[:50] + "..." if len(idea) > 50 else idea
            branch = {
                'name': f"فكرة {i+1}",
                'children': [{'name': short_idea}]
            }
            mindmap_data['branches'].append(branch)
        
        # الحد من عدد الفروع
        if len(mindmap_data['branches']) > 5:
            mindmap_data['branches'] = mindmap_data['branches'][:5]
        
        return mindmap_data

